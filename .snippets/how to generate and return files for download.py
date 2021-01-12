from flask import Flask, send_file
from pptx import Presentation

app = Flask(__name__)

@app.route('/')
def index():
    return '<body>\
                <a href=/generate-file/>\
                        <button class="btn btn-default">\
                            Generate Powerpoint\
                        </button>\
                    </a>\
                    \
                    <a href=/return-file/>\
                        <button class="btn btn-default">\
                            Download\
                        </button>\
                    </a>\
                </body>'

@app.route('/generate-file/')
def generate_file():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Okay!"
    subtitle.text = "I figured it out!"

    prs.save('test.pptx')
    return 'Hello'

@app.route('/return-file/')
def return_file():
    return send_file('test.pptx', as_attachment=True,attachment_filename='test.pptx', cache_timeout=0)

if __name__ == ('__main__'):
    app.run(debug=True)
