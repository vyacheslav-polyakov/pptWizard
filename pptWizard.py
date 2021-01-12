from pptx import Presentation
from pptx.util import Inches, Pt

from txtProcessor import formatTitle, Limit, findKeywords

# Crawl the web and get a presentation structure
import crawler
question = "What's the topic of your presentation: "
while True:
    structure = Limit(crawler.Search(input(question)))
    if structure == None:
        question = "Try to rephrase your request: "
        continue
    structure.append(crawler.attachImages(structure))
    break

# Extract items from the structure
main_title = structure[0]
pptname = formatTitle(main_title)
headlines = structure[1]
paragraphs = structure[2]
images = structure[3]

# Define inches and create a blank presentation
inch = Inches(1.0)
ppt = Presentation()

# Add a title slide
slide = ppt.slides.add_slide(ppt.slide_layouts[6])
# Add the main title
titlebox = slide.shapes.add_textbox(
    left = inch,
    top = inch*2,
    width = inch*8,
    height = inch*3
    )
frame = titlebox.text_frame
frame.word_wrap = True
title = frame.add_paragraph()
title.text = main_title
title.font.bold = True
title.font.size = Pt(50)

# Create slides from all paragraphs
for paragraph in paragraphs:

    # Set default values
    textWidth = inch*5.0
    textTop = inch*1.5

    # Create a blank slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    # See if there is an image to attach:
    path = images[paragraph]
    if path != None:
        picture = slide.shapes.add_picture(
                path,
                left = inch*5.5,
                top = inch*1.5,
                width = inch*4,
                height = inch*5
                )
    # If there's no picture
    else:
        # Allign the default parameters of the text to put it in the middle
        textWidth *= 2
        textTop += 0.2*inch
        pass

    # See if there is a headline to add
    headline = headlines[paragraphs.index(paragraph)]

    # The shorter the headine is the more it should shift to the middle
    coef = 3
    length = len(headline)
    while coef > 1 and length > 10:
        coef -= 1
        length -= 10
    headLeft = inch*coef

    if headline != None:
        # Add the title
        titlebox = slide.shapes.add_textbox(
            left = headLeft,
            top = 0,
            width = inch*9,
            height = inch
            )
        frame = titlebox.text_frame
        frame.word_wrap = True
        title = frame.add_paragraph()
        title.text = headline
        title.font.bold = True
        title.font.size = Pt(40)
    else:
        pass

    # Add the paragraph
    textbox = slide.shapes.add_textbox(
        left = inch*0.2,
        top = textTop,
        width = textWidth,
        height = inch*4
        )
    frame = textbox.text_frame
    frame.word_wrap = True
    content = frame.add_paragraph()
    content.text = paragraph
    content.font.size = Pt(20)

# Save the presentation file
ppt.save("{}.pptx".format(pptname))
print("\npptWizard program complete")
input()
